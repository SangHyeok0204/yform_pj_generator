"""Y-PJ-FoRM PPTX rendering engine.

Encodes prompts/design_system.md as Python constants and slide builders.
Content-agnostic: takes a DeckContent object, returns PPTX bytes.
"""
from io import BytesIO
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

from .schema import DeckContent

# === Paths ===
PROJECT_ROOT = Path(__file__).resolve().parent.parent.parent
ASSETS_DIR = PROJECT_ROOT / "assets" / "images"

# === Colors ===
NAVY = RGBColor(0x12, 0x2B, 0x46)
INK = RGBColor(0x00, 0x00, 0x00)
RED = RGBColor(0xC0, 0x00, 0x00)
INK_M = RGBColor(0x7A, 0x7A, 0x7A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
HAIR = RGBColor(0xE0, 0xE0, 0xE0)
SOFT = RGBColor(0xF0, 0xF0, 0xF0)

# === Fonts ===
F_LIGHT = "Pretendard ExtraLight"
F_SEMI = "Pretendard SemiBold"
F_BOLD = "Pretendard ExtraBold"
F_REG = "Pretendard"

# === Cover anchors === (matched to reference/master_deck_1.pptx slide 1)
COVER_HERO_MAIN = (Inches(0.00), Inches(0.00), Inches(10.09), Inches(6.70))
COVER_HERO_STRIP = (Inches(10.09), Inches(6.70), Inches(3.24), Inches(0.80))
COVER_LOGO = (Inches(10.58), Inches(-0.42), Inches(2.27), Inches(1.70))
COVER_TITLE = (Inches(0.27), Inches(2.91), Inches(9.40), Inches(2.42))
COVER_MEMBERS = (Inches(0.20), Inches(6.88), Inches(9.41), Inches(0.44))

# === Standard anchors ===
TITLE_X, TITLE_Y, TITLE_W, TITLE_H = Inches(0.12), Inches(0.09), Inches(13.09), Inches(0.77)
RULE_X1, RULE_Y, RULE_X2 = Inches(0.00), Inches(0.93), Inches(13.33)
SUB_X, SUB_Y, SUB_W, SUB_H = Inches(0.14), Inches(1.00), Inches(13.05), Inches(0.55)
BODY_X, BODY_Y, BODY_W, BODY_H = Inches(0.12), Inches(1.62), Inches(13.09), Inches(5.58)

# === TOC anchors === (matched to reference/master_deck_1.pptx slide 2)
TOC_DECOR = (Inches(10.12), Inches(-0.04), Inches(11.29), Inches(7.50))
TOC_HDR_X, TOC_HDR_Y, TOC_HDR_W, TOC_HDR_H = Inches(0.31), Inches(0.11), Inches(4.51), Inches(0.94)
TOC_RULE_X1, TOC_RULE_Y, TOC_RULE_X2 = Inches(0.00), Inches(1.20), Inches(10.12)
TOC_LIST_X, TOC_LIST_Y, TOC_LIST_W, TOC_LIST_H = Inches(0.31), Inches(1.34), Inches(9.45), Inches(3.13)
TOC_BOTTOM_RULE_Y = Inches(6.70)


# === Helpers ===
def add_textbox(slide, left, top, width, height):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tb.text_frame.word_wrap = True
    tb.text_frame.margin_left = 0
    tb.text_frame.margin_right = 0
    tb.text_frame.margin_top = 0
    tb.text_frame.margin_bottom = 0
    return tb


def style_run(run, font=F_LIGHT, size=20, color=INK, bold=False, spacing=None):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = etree.SubElement(rPr, qn("a:ea"))
    ea.set("typeface", font)
    cs = rPr.find(qn("a:cs"))
    if cs is None:
        cs = etree.SubElement(rPr, qn("a:cs"))
    cs.set("typeface", font)
    if spacing is not None:
        rPr.set("spc", str(int(spacing)))


def add_hairline(slide, x1=RULE_X1, y=RULE_Y, x2=RULE_X2, color=NAVY, weight=0.75):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y, x2, y)
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    return line


def add_section_title(slide, text):
    tb = add_textbox(slide, TITLE_X, TITLE_Y, TITLE_W, TITLE_H)
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = text
    style_run(r, font=F_SEMI, size=40, color=NAVY, bold=True, spacing=-40)
    return tb


def add_subtitle(slide, text):
    tb = add_textbox(slide, SUB_X, SUB_Y, SUB_W, SUB_H)
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = text
    style_run(r, font=F_SEMI, size=20, color=NAVY, bold=True)
    return tb


def add_body_paragraphs(
    slide,
    blocks,
    x=BODY_X,
    y=BODY_Y,
    w=BODY_W,
    h=BODY_H,
    body_size=18,
    head_size=20,
    lead_space=10,
    line_space=1.18,
):
    tb = add_textbox(slide, x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for blk in blocks:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        kind = blk.get("kind", "p")
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_space
        if kind == "h":
            p.space_before = Pt(blk.get("space_before", lead_space))
            p.space_after = Pt(blk.get("space_after", 2))
            r = p.add_run()
            r.text = blk["text"]
            style_run(r, font=F_SEMI, size=blk.get("size") or head_size, color=NAVY, bold=True)
        elif kind == "p":
            p.space_before = Pt(blk.get("space_before", 2))
            p.space_after = Pt(blk.get("space_after", 2))
            r = p.add_run()
            r.text = blk["text"]
            style_run(r, font=F_LIGHT, size=blk.get("size") or body_size, color=INK)
        elif kind == "li":
            p.space_before = Pt(blk.get("space_before", 4))
            p.space_after = Pt(blk.get("space_after", 2))
            num = blk.get("num") or ""
            label = blk.get("label") or ""
            text = blk.get("text", "")
            size = blk.get("size") or body_size
            if num:
                r0 = p.add_run()
                r0.text = num + " "
                style_run(r0, font=F_SEMI, size=size, color=NAVY, bold=True)
            if label:
                r1 = p.add_run()
                r1.text = label + " "
                style_run(r1, font=F_SEMI, size=size, color=NAVY, bold=True)
                rdash = p.add_run()
                rdash.text = "— "
                style_run(rdash, font=F_LIGHT, size=size, color=INK_M)
            r2 = p.add_run()
            r2.text = text
            style_run(r2, font=F_LIGHT, size=size, color=INK)
        elif kind == "gap":
            p.space_after = Pt(blk.get("pt", 6))
            r = p.add_run()
            r.text = ""
            style_run(r, font=F_LIGHT, size=4, color=INK)
        elif kind == "caption":
            p.space_before = Pt(blk.get("space_before", 4))
            p.space_after = Pt(blk.get("space_after", 0))
            r = p.add_run()
            r.text = blk["text"]
            style_run(r, font=F_REG, size=blk.get("size") or 12, color=INK_M)
    return tb


# === Slide builders ===
def build_cover(prs, project_title, project_subtitle, member_line):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(str(ASSETS_DIR / "cover_hero_main.png"), *COVER_HERO_MAIN)
    slide.shapes.add_picture(str(ASSETS_DIR / "cover_hero_strip.png"), *COVER_HERO_STRIP)
    slide.shapes.add_picture(str(ASSETS_DIR / "logo.jpg"), *COVER_LOGO)

    tb = add_textbox(slide, *COVER_TITLE)
    tf = tb.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.line_spacing = 1.05
    r = p.add_run()
    r.text = project_title
    style_run(r, font=F_BOLD, size=40, color=WHITE, bold=True, spacing=-80)

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.LEFT
    p2.line_spacing = 1.15
    p2.space_before = Pt(10)
    r2 = p2.add_run()
    r2.text = project_subtitle
    style_run(r2, font=F_SEMI, size=18, color=WHITE, bold=True)

    tb2 = add_textbox(slide, *COVER_MEMBERS)
    p = tb2.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = member_line
    style_run(r, font=F_SEMI, size=20, color=INK, bold=True)
    return slide


def build_toc(prs, items):
    """Build TOC verbatim per master_deck_1.pptx slide 2:
    decor strip (right) -> 'Contents' SemiBold 50pt -> short top rule ->
    single-textbox numbered list (one paragraph per section, ExtraLight 20pt navy,
    'N. Name' as one run -- no red, no chip) -> full-width bottom rule.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(str(ASSETS_DIR / "toc_decor.png"), *TOC_DECOR)

    tb = add_textbox(slide, TOC_HDR_X, TOC_HDR_Y, TOC_HDR_W, TOC_HDR_H)
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = "Contents"
    style_run(r, font=F_SEMI, size=50, color=NAVY, bold=True, spacing=-100)

    add_hairline(slide, TOC_RULE_X1, TOC_RULE_Y, TOC_RULE_X2, color=NAVY)

    tb = add_textbox(slide, TOC_LIST_X, TOC_LIST_Y, TOC_LIST_W, TOC_LIST_H)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for i, item in enumerate(items, 1):
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.5
        p.space_before = Pt(2)
        p.space_after = Pt(2)
        r = p.add_run()
        r.text = f"{i}. {item}"
        style_run(r, font=F_LIGHT, size=20, color=NAVY)

    add_hairline(slide, Inches(0.00), TOC_BOTTOM_RULE_Y, Inches(13.33), color=NAVY)
    return slide


def build_content(prs, title, subtitle, blocks, body_size=18, head_size=20):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_title(slide, title)
    add_hairline(slide)
    if subtitle:
        add_subtitle(slide, subtitle)
        add_body_paragraphs(slide, blocks, body_size=body_size, head_size=head_size)
    else:
        add_body_paragraphs(
            slide,
            blocks,
            x=BODY_X,
            y=Inches(1.05),
            w=BODY_W,
            h=Inches(6.15),
            body_size=body_size,
            head_size=head_size,
        )
    return slide


# === Public entry ===
def render_deck(content: DeckContent) -> bytes:
    """Render a DeckContent into a PPTX file as bytes."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    build_cover(
        prs,
        project_title=content.cover.title,
        project_subtitle=content.cover.subtitle,
        member_line=content.cover.member_line,
    )
    build_toc(prs, content.toc)
    for slide in content.slides:
        blocks = [b.model_dump() for b in slide.blocks]
        build_content(prs, slide.title, slide.subtitle, blocks)

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()
