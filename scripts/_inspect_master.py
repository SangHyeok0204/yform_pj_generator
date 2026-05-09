"""Inspect master_deck_1.pptx and master_deck_2.pptx structure: shapes, fonts, sizes, positions."""
import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ROOT))

from pptx import Presentation
from pptx.util import Emu

DECKS = [
    Path(r"C:\temp_decks\m1.pptx"),
    Path(r"C:\temp_decks\m2.pptx"),
]
OUT = ROOT / "scripts" / "_master_dump.txt"


def emu_to_in(v):
    return v / 914400 if v is not None else None


def inspect(path: Path, max_slides: int, lines: list[str]):
    with open(path, "rb") as fh:
        prs = Presentation(fh)
    lines.append(f"\n========== {path.name} ==========")
    lines.append(f"slide_width:  {emu_to_in(prs.slide_width):.3f}\"")
    lines.append(f"slide_height: {emu_to_in(prs.slide_height):.3f}\"")
    lines.append(f"slides:       {len(prs.slides)}")
    for si, slide in enumerate(prs.slides):
        if si >= max_slides:
            break
        lines.append(f"\n--- slide {si + 1} ({len(slide.shapes)} shapes) ---")
        for sh in slide.shapes:
            try:
                left = emu_to_in(sh.left)
                top = emu_to_in(sh.top)
                w = emu_to_in(sh.width)
                h = emu_to_in(sh.height)
            except Exception:
                left = top = w = h = None
            kind = sh.shape_type
            name = sh.name
            lines.append(
                f"  [{kind}] {name}  L={left} T={top} W={w} H={h}"
            )
            if sh.has_text_frame:
                for pi, p in enumerate(sh.text_frame.paragraphs):
                    align = p.alignment
                    for ri, r in enumerate(p.runs):
                        f = r.font
                        col = None
                        try:
                            col = f.color.rgb
                        except Exception:
                            pass
                        lines.append(
                            f"    p{pi}r{ri}: align={align} "
                            f"font={f.name!r} size={f.size} bold={f.bold} "
                            f"color={col!r} text={r.text!r}"
                        )


def main():
    lines: list[str] = []
    for d in DECKS:
        if not d.exists():
            lines.append(f"MISSING: {d}")
            continue
        inspect(d, max_slides=6, lines=lines)
    OUT.write_text("\n".join(lines), encoding="utf-8")
    print(f"wrote {len(lines)} lines -> {OUT.relative_to(ROOT)}")


if __name__ == "__main__":
    main()
