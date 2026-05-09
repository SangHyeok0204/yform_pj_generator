"""Verify generated deck matches master_deck_1.pptx structural anchors.

Compares cover/TOC anchors and counts shapes per slide. Cross-checks against
reference/master_deck_1.pptx (slides 1-3) and design_system.md hard rules.
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor

OUT = Path(r"C:\temp_decks\out.pptx")
MASTER = Path(r"C:\temp_decks\m1.pptx")
DUMP = Path(r"C:\temp_decks\verify.txt")
_LINES: list[str] = []


def say(*args):
    line = " ".join(str(a) for a in args)
    _LINES.append(line)


def flush():
    DUMP.write_text("\n".join(_LINES), encoding="utf-8")


def emu_to_in(v):
    return v / 914400 if v is not None else None


def almost(a, b, eps=0.05):
    if a is None or b is None:
        return False
    return abs(a - b) <= eps


def shape_summary(slide):
    out = []
    for sh in slide.shapes:
        try:
            l = round(emu_to_in(sh.left), 2)
            t = round(emu_to_in(sh.top), 2)
            w = round(emu_to_in(sh.width), 2)
            h = round(emu_to_in(sh.height), 2)
        except Exception:
            l = t = w = h = None
        out.append((sh.shape_type, sh.name, l, t, w, h))
    return out


def main():
    out_prs = Presentation(open(OUT, "rb"))
    master_prs = Presentation(open(MASTER, "rb"))

    say(f"=== generated deck ({len(out_prs.slides)} slides) ===\n")

    # ---- Cover (slide 1) ----
    say("--- Cover (slide 1) ---")
    cs = out_prs.slides[0]
    for sh in cs.shapes:
        try:
            l = round(emu_to_in(sh.left), 3)
            t = round(emu_to_in(sh.top), 3)
            w = round(emu_to_in(sh.width), 3)
            h = round(emu_to_in(sh.height), 3)
        except Exception:
            l = t = w = h = None
        say(f"  [{sh.shape_type}] {sh.name} L={l} T={t} W={w} H={h}")
        if sh.has_text_frame:
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    sz = r.font.size.pt if r.font.size else None
                    say(f"     -> font={r.font.name!r} size={sz}pt bold={r.font.bold} text={r.text!r}")

    # ---- TOC (slide 2) ----
    say("\n--- TOC (slide 2) ---")
    ts = out_prs.slides[1]
    for sh in ts.shapes:
        try:
            l = round(emu_to_in(sh.left), 3)
            t = round(emu_to_in(sh.top), 3)
            w = round(emu_to_in(sh.width), 3)
            h = round(emu_to_in(sh.height), 3)
        except Exception:
            l = t = w = h = None
        say(f"  [{sh.shape_type}] {sh.name} L={l} T={t} W={w} H={h}")
        if sh.has_text_frame:
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    sz = r.font.size.pt if r.font.size else None
                    col = None
                    try:
                        col = r.font.color.rgb
                    except Exception:
                        pass
                    say(
                        f"     -> font={r.font.name!r} size={sz}pt bold={r.font.bold} "
                        f"color={col} text={r.text!r}"
                    )

    # ---- Verifications ----
    say("\n=== HARD-RULE CHECKS ===\n")
    fails = []

    # 1. Slide canvas
    if abs(emu_to_in(out_prs.slide_width) - 13.333) > 0.01:
        fails.append(f"slide_width != 13.333 ({emu_to_in(out_prs.slide_width)})")
    if abs(emu_to_in(out_prs.slide_height) - 7.5) > 0.01:
        fails.append(f"slide_height != 7.5 ({emu_to_in(out_prs.slide_height)})")

    # 2. Cover slide must have logo, hero_main, hero_strip pictures
    cs_pics = [sh for sh in cs.shapes if sh.shape_type == 13]
    if len(cs_pics) < 3:
        fails.append(f"cover has only {len(cs_pics)} pictures (need 3: hero_main, hero_strip, logo)")

    # 3. TOC must have decor picture
    ts_pics = [sh for sh in ts.shapes if sh.shape_type == 13]
    if len(ts_pics) < 1:
        fails.append("TOC has no decor picture (toc_decor.png missing)")
    else:
        decor = ts_pics[0]
        decor_l = round(emu_to_in(decor.left), 2)
        decor_t = round(emu_to_in(decor.top), 2)
        decor_w = round(emu_to_in(decor.width), 2)
        if not (almost(decor_l, 10.12, 0.05) and almost(decor_w, 11.29, 0.1)):
            fails.append(
                f"TOC decor anchor wrong (L={decor_l} W={decor_w}, expected L≈10.12 W≈11.29)"
            )

    # 4. TOC must have NO red text and "Contents" must be SemiBold (not ExtraBold)
    contents_found = False
    for sh in ts.shapes:
        if sh.has_text_frame:
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    if r.text == "Contents":
                        contents_found = True
                        if r.font.name != "Pretendard SemiBold":
                            fails.append(
                                f"TOC 'Contents' is {r.font.name!r}, must be 'Pretendard SemiBold'"
                            )
                        if r.font.size and r.font.size.pt != 50:
                            fails.append(
                                f"TOC 'Contents' size is {r.font.size.pt}pt, must be 50pt"
                            )
                    try:
                        col = r.font.color.rgb
                        if col == RGBColor(0xC0, 0x00, 0x00):
                            fails.append(f"TOC has RED text: {r.text!r}")
                    except Exception:
                        pass
    if not contents_found:
        fails.append("TOC missing 'Contents' header")

    # 5. TOC must have 2 lines (top rule + bottom rule)
    ts_lines = [sh for sh in ts.shapes if sh.shape_type == 9]
    if len(ts_lines) < 2:
        fails.append(f"TOC has only {len(ts_lines)} hairlines (need 2: top + bottom)")

    # 6. Every non-cover, non-TOC slide must have a hairline at Y=0.93
    for i, slide in enumerate(out_prs.slides):
        if i < 2:
            continue
        rule_ys = []
        for sh in slide.shapes:
            if sh.shape_type == 9:
                ty = round(emu_to_in(sh.top), 2)
                rule_ys.append(ty)
        if not any(almost(y, 0.93, 0.05) for y in rule_ys):
            fails.append(f"slide {i + 1} missing hairline at Y≈0.93 (found rules at {rule_ys})")

    # 7. Section title must be present on every non-cover, non-TOC slide
    for i, slide in enumerate(out_prs.slides):
        if i < 2:
            continue
        has_title = False
        for sh in slide.shapes:
            if sh.has_text_frame:
                for p in sh.text_frame.paragraphs:
                    for r in p.runs:
                        if r.font.size and r.font.size.pt == 40:
                            has_title = True
        if not has_title:
            fails.append(f"slide {i + 1} missing 40pt section title")

    # 8. No logo on non-cover slides — check by detecting a 2.27"x1.70" picture
    for i, slide in enumerate(out_prs.slides):
        if i == 0:
            continue
        for sh in slide.shapes:
            if sh.shape_type == 13:
                w = round(emu_to_in(sh.width), 2)
                h = round(emu_to_in(sh.height), 2)
                if almost(w, 2.27, 0.05) and almost(h, 1.70, 0.05):
                    fails.append(f"slide {i + 1} has logo-sized picture (forbidden)")

    # 9. Master comparison — TOC contents header position
    master_ts = master_prs.slides[1]
    master_contents = None
    for sh in master_ts.shapes:
        if sh.has_text_frame:
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    if r.text == "Contents":
                        master_contents = sh
    out_contents = None
    for sh in ts.shapes:
        if sh.has_text_frame:
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    if r.text == "Contents":
                        out_contents = sh
    if master_contents and out_contents:
        for attr, name in [("left", "L"), ("top", "T"), ("width", "W"), ("height", "H")]:
            mv = round(emu_to_in(getattr(master_contents, attr)), 2)
            ov = round(emu_to_in(getattr(out_contents, attr)), 2)
            if not almost(mv, ov, 0.05):
                fails.append(f"TOC 'Contents' {name}: master={mv} vs ours={ov}")

    # Report
    if fails:
        say(f"FAIL — {len(fails)} issues:")
        for f in fails:
            say(f"  - {f}")
        return 1
    else:
        say("ALL CHECKS PASSED")
        say(f"  - cover pictures: {len(cs_pics)}")
        say(f"  - TOC pictures: {len(ts_pics)}, lines: {len(ts_lines)}")
        say(f"  - non-cover/TOC slides: {len(out_prs.slides) - 2}")
        return 0


if __name__ == "__main__":
    rc = main()
    flush()
    print(f"verify rc={rc}; details -> {DUMP}")
    raise SystemExit(rc)
