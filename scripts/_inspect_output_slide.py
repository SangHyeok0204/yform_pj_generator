"""Dump shapes for slides 3 and 6 of the generated deck."""
from pathlib import Path

from pptx import Presentation

OUT = Path(r"C:\temp_decks\out.pptx")
DUMP = Path(r"C:\temp_decks\out_slides.txt")


def emu_to_in(v):
    return v / 914400 if v is not None else None


def main():
    prs = Presentation(open(OUT, "rb"))
    lines: list[str] = []
    for idx in [2, 5, 13, 27]:
        slide = prs.slides[idx]
        lines.append(f"\n=== slide {idx + 1} ({len(slide.shapes)} shapes) ===")
        for sh in slide.shapes:
            try:
                l = round(emu_to_in(sh.left), 3)
                t = round(emu_to_in(sh.top), 3)
                w = round(emu_to_in(sh.width), 3)
                h = round(emu_to_in(sh.height), 3)
            except Exception:
                l = t = w = h = None
            lines.append(f"  [{sh.shape_type}] {sh.name} L={l} T={t} W={w} H={h}")
            if sh.has_text_frame:
                for pi, p in enumerate(sh.text_frame.paragraphs):
                    bits = []
                    for r in p.runs:
                        sz = r.font.size.pt if r.font.size else None
                        bits.append(f"{r.font.name!r}/{sz}/{r.font.bold} {r.text!r}")
                    if bits:
                        lines.append(f"    p{pi}: {bits[0]}" + (" ..." if len(bits) > 1 else ""))
    DUMP.write_text("\n".join(lines), encoding="utf-8")
    print(f"-> {DUMP}")


if __name__ == "__main__":
    main()
