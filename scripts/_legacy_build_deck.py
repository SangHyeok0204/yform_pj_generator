# -*- coding: utf-8 -*-
"""
Build a Y-PJ-FoRM academic-society deck for the paper:
"Options Market Makers" (Hu, Kirilova, Muravyev, Ryu, 2025)
KOSPI 200 Delta Hedging study.

Follows template.md anchors, fonts, hairlines, and density rules.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from copy import deepcopy
from lxml import etree

ROOT = os.path.dirname(os.path.abspath(__file__))

# ---- Colors ----
NAVY  = RGBColor(0x12, 0x2B, 0x46)
INK   = RGBColor(0x00, 0x00, 0x00)
RED   = RGBColor(0xC0, 0x00, 0x00)
INK_M = RGBColor(0x7A, 0x7A, 0x7A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
HAIR  = RGBColor(0xE0, 0xE0, 0xE0)
SOFT  = RGBColor(0xF0, 0xF0, 0xF0)

# ---- Fonts ----
F_LIGHT = "Pretendard ExtraLight"
F_SEMI  = "Pretendard SemiBold"
F_BOLD  = "Pretendard ExtraBold"
F_REG   = "Pretendard"

# ---- Cover anchors ----
COVER_HERO_MAIN  = (Inches(0.00),  Inches(0.00), Inches(10.09), Inches(6.70))
COVER_HERO_STRIP = (Inches(10.09), Inches(6.70), Inches(3.24),  Inches(0.80))
COVER_LOGO       = (Inches(10.58), Inches(1.14), Inches(2.27),  Inches(1.70))
COVER_TITLE      = (Inches(0.50),  Inches(2.91), Inches(9.20),  Inches(2.42))
COVER_MEMBERS    = (Inches(0.20),  Inches(6.88), Inches(9.41),  Inches(0.44))

# ---- Standard anchors ----
TITLE_X, TITLE_Y, TITLE_W, TITLE_H = Inches(0.12), Inches(0.09), Inches(13.09), Inches(0.77)
RULE_X1, RULE_Y, RULE_X2 = Inches(0.00), Inches(0.93), Inches(13.33)
SUB_X, SUB_Y, SUB_W, SUB_H = Inches(0.14), Inches(1.00), Inches(13.05), Inches(0.55)
BODY_X, BODY_Y, BODY_W, BODY_H = Inches(0.12), Inches(1.62), Inches(13.09), Inches(5.58)

# ---- TOC anchors ----
TOC_HDR_X, TOC_HDR_Y, TOC_HDR_W, TOC_HDR_H = Inches(0.31), Inches(0.11), Inches(4.51), Inches(0.94)
TOC_RULE_X1, TOC_RULE_Y, TOC_RULE_X2 = Inches(0.00), Inches(1.20), Inches(10.12)
TOC_LIST_X, TOC_LIST_Y, TOC_LIST_W, TOC_LIST_H = Inches(0.31), Inches(1.34), Inches(9.45), Inches(4.75)
TOC_BOTTOM_RULE_Y = Inches(6.70)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------
def add_textbox(slide, left, top, width, height):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tb.text_frame.word_wrap = True
    tb.text_frame.margin_left = 0
    tb.text_frame.margin_right = 0
    tb.text_frame.margin_top = 0
    tb.text_frame.margin_bottom = 0
    return tb

def style_run(run, font=F_LIGHT, size=20, color=INK, bold=False, spacing=None):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    # apply East Asian font as well so Korean uses Pretendard
    rPr = run._r.get_or_add_rPr()
    # set ea (east asian) font
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = etree.SubElement(rPr, qn('a:ea'))
    ea.set('typeface', font)
    cs = rPr.find(qn('a:cs'))
    if cs is None:
        cs = etree.SubElement(rPr, qn('a:cs'))
    cs.set('typeface', font)
    if spacing is not None:
        # spacing is in 1/100 of a point; -2% of size translates to -2 * size/100 hundredths
        rPr.set('spc', str(int(spacing)))
    # Korean line break behavior is enabled at paragraph level via word_wrap

def set_para_alignment(p, align):
    p.alignment = align

def set_para_linebreak_keepall(p):
    pPr = p._pPr if p._pPr is not None else p._p.get_or_add_pPr()
    # eastAsianLineBreak / wordBreak: keep-all approximated via kumimoji not applicable.
    # python-pptx limit; rely on word_wrap and design discipline.
    pass

def add_hairline(slide, x1=RULE_X1, y=RULE_Y, x2=RULE_X2, color=NAVY, weight=0.75):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y, x2, y)
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    return line

def add_section_title(slide, text):
    tb = add_textbox(slide, TITLE_X, TITLE_Y, TITLE_W, TITLE_H)
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = text
    style_run(r, font=F_SEMI, size=40, color=NAVY, bold=True, spacing=-40)
    return tb

def add_subtitle(slide, text):
    tb = add_textbox(slide, SUB_X, SUB_Y, SUB_W, SUB_H)
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = text
    style_run(r, font=F_SEMI, size=24, color=NAVY, bold=True)
    return tb

def add_body_paragraphs(slide, blocks, x=BODY_X, y=BODY_Y, w=BODY_W, h=BODY_H,
                       body_size=18, head_size=20, lead_space=10, line_space=1.18):
    """blocks: list of dicts:
       {'kind':'h', 'text':'...'}        SemiBold-navy sub-head
       {'kind':'p', 'text':'...'}        body paragraph
       {'kind':'li', 'num':'1.', 'label':'...', 'text':'...'}  numbered list
       {'kind':'gap', 'pt':6}            extra space
    """
    tb = add_textbox(slide, x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for blk in blocks:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        kind = blk.get('kind', 'p')
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_space
        if kind == 'h':
            p.space_before = Pt(blk.get('space_before', lead_space))
            p.space_after = Pt(blk.get('space_after', 2))
            r = p.add_run()
            r.text = blk['text']
            style_run(r, font=F_SEMI, size=blk.get('size', head_size),
                      color=NAVY, bold=True)
        elif kind == 'p':
            p.space_before = Pt(blk.get('space_before', 2))
            p.space_after = Pt(blk.get('space_after', 2))
            r = p.add_run()
            r.text = blk['text']
            style_run(r, font=F_LIGHT, size=blk.get('size', body_size), color=INK)
        elif kind == 'li':
            p.space_before = Pt(blk.get('space_before', 4))
            p.space_after = Pt(blk.get('space_after', 2))
            num = blk.get('num', '')
            label = blk.get('label', '')
            text = blk.get('text', '')
            if num:
                r0 = p.add_run()
                r0.text = num + ' '
                style_run(r0, font=F_SEMI, size=blk.get('size', body_size),
                          color=NAVY, bold=True)
            if label:
                r1 = p.add_run()
                r1.text = label + ' '
                style_run(r1, font=F_SEMI, size=blk.get('size', body_size),
                          color=NAVY, bold=True)
                rdash = p.add_run()
                rdash.text = '— '
                style_run(rdash, font=F_LIGHT, size=blk.get('size', body_size),
                          color=INK_M)
            r2 = p.add_run()
            r2.text = text
            style_run(r2, font=F_LIGHT, size=blk.get('size', body_size), color=INK)
        elif kind == 'gap':
            p.space_after = Pt(blk.get('pt', 6))
            r = p.add_run()
            r.text = ''
            style_run(r, font=F_LIGHT, size=4, color=INK)
        elif kind == 'caption':
            p.space_before = Pt(blk.get('space_before', 4))
            p.space_after = Pt(blk.get('space_after', 0))
            r = p.add_run()
            r.text = blk['text']
            style_run(r, font=F_REG, size=blk.get('size', 12), color=INK_M)
    return tb


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------
def build_cover(prs, project_title, project_subtitle, member_line):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # Hero images
    slide.shapes.add_picture(os.path.join(ROOT, 'cover_hero_main.png'),
                             *COVER_HERO_MAIN)
    slide.shapes.add_picture(os.path.join(ROOT, 'cover_hero_strip.png'),
                             *COVER_HERO_STRIP)
    # Logo
    slide.shapes.add_picture(os.path.join(ROOT, 'logo.jpg'),
                             *COVER_LOGO)
    # Title
    tb = add_textbox(slide, *COVER_TITLE)
    tf = tb.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.line_spacing = 1.05
    r = p.add_run()
    r.text = project_title
    style_run(r, font=F_BOLD, size=40, color=WHITE, bold=True, spacing=-80)

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.LEFT
    p2.line_spacing = 1.15
    p2.space_before = Pt(10)
    r2 = p2.add_run()
    r2.text = project_subtitle
    style_run(r2, font=F_SEMI, size=18, color=WHITE, bold=True)

    # Member roster
    tb2 = add_textbox(slide, *COVER_MEMBERS)
    p = tb2.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = member_line
    style_run(r, font=F_SEMI, size=20, color=INK, bold=True)
    return slide


def build_toc(prs, items):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Header
    tb = add_textbox(slide, TOC_HDR_X, TOC_HDR_Y, TOC_HDR_W, TOC_HDR_H)
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = 'Contents'
    style_run(r, font=F_BOLD, size=50, color=NAVY, bold=True, spacing=-100)

    add_hairline(slide, TOC_RULE_X1, TOC_RULE_Y, TOC_RULE_X2, color=NAVY)

    # List
    tb = add_textbox(slide, TOC_LIST_X, TOC_LIST_Y, TOC_LIST_W, TOC_LIST_H)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for i, item in enumerate(items, 1):
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.5
        p.space_before = Pt(6)
        p.space_after = Pt(6)
        rn = p.add_run()
        rn.text = f"{i:02d}.   "
        style_run(rn, font=F_BOLD, size=20, color=RED, bold=True)
        rt = p.add_run()
        rt.text = item
        style_run(rt, font=F_LIGHT, size=20, color=INK)

    # Optional bottom rule
    add_hairline(slide, Inches(0.00), TOC_BOTTOM_RULE_Y, Inches(13.33), color=NAVY)
    return slide


def build_content(prs, title, subtitle, blocks, body_size=18, head_size=20):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_title(slide, title)
    add_hairline(slide)
    if subtitle:
        add_subtitle(slide, subtitle)
        add_body_paragraphs(slide, blocks, body_size=body_size, head_size=head_size)
    else:
        # subtitle omitted, body shifts up to Y=1.05, height 6.15"
        add_body_paragraphs(slide, blocks,
                            x=BODY_X, y=Inches(1.05),
                            w=BODY_W, h=Inches(6.15),
                            body_size=body_size, head_size=head_size)
    return slide


# ---------------------------------------------------------------------------
# Content
# ---------------------------------------------------------------------------
TOC_ITEMS = [
    "Research Question & Motivation",
    "Data & Market Maker Identification",
    "Profitability of OMMs",
    "Delta Hedging Behavior",
    "Active Inventory Management",
    "Evidence from S&P 500 Options",
    "Conclusion & Implications",
]

def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ---- Cover ----
    build_cover(
        prs,
        project_title="Options Market Makers",
        project_subtitle="How OMMs Trade and Manage Risk in KOSPI 200 Index Options",
        member_line="Hu · Kirilova · Muravyev · Ryu  (2025)   |   Reviewed by Y-FoRM",
    )

    # ---- TOC ----
    build_toc(prs, TOC_ITEMS)

    # =========================================================================
    # 1. Research Question & Motivation
    # =========================================================================
    build_content(
        prs,
        title="1. Research Question & Motivation",
        subtitle="1.1 Microstructure 이론과 옵션 이론의 충돌",
        blocks=[
            {'kind':'p','text':"헷지(hedging)는 위험관리의 핵심 전략이다. 옵션의 방향성 위험은 기초자산을 거래하는 델타헷지(delta hedging)로 상쇄할 수 있다고 알려져 왔다. 그러나 마이크로스트럭처 이론과 옵션 이론은 시장조성자의 위험관리에 대해 서로 다른 관점을 제시한다."},

            {'kind':'h','text':"마이크로스트럭처 이론의 시각"},
            {'kind':'p','text':"Stoll(1978), Ho and Stoll(1983) 등은 헷지보다 재고관리(inventory management)와 역선택위험(adverse selection)에 주목한다. 시장조성자는 양방향 호가를 조정해 균형 잡힌 주문 흐름을 유도하고, 재고가 한쪽으로 쏠리면 호가를 비대칭적으로 움직여 자연스럽게 평형을 회복시킨다."},

            {'kind':'h','text':"옵션 이론과 교과서의 시각"},
            {'kind':'p','text':"Black–Scholes(1972), Leland(1985), Hull(2018) 등 옵션 문헌은 델타헷지를 핵심 위험관리 수단으로 본다. Hull and White(2017)는 \"옵션 트레이더는 빈번하게 델타를 조정해 0에 가깝게 유지한다\"고 단정한다. 즉, 시장조성자가 항상 풀(full) 델타헷지를 한다는 것이 통념이다."},

            {'kind':'h','text':"본 연구의 핵심 질문"},
            {'kind':'p','text':"마이크로스트럭처 이론은 위험관리의 핵심을 놓쳤는가, 아니면 옵션 이론의 통념이 과장되었는가? 실제 옵션 시장조성자(OMM)는 어떤 방식으로 수익을 내고 위험을 관리하는가? 이 질문은 직접적 거래 데이터 없이는 답할 수 없다."},
        ],
    )

    build_content(
        prs,
        title="1. Research Question & Motivation",
        subtitle="1.2 본 연구의 기여 — OMM에 대한 최초의 직접 분석",
        blocks=[
            {'kind':'p','text':"전 세계에서 2023년 한 해에만 1,080억 건의 옵션 계약이 거래되었다. 미국에서만 135만 개의 서로 다른 옵션 종목이 상장되어 있고, 시장조성자(OMM)는 이들 대부분의 두 방향(two-sided) 호가를 사실상 단독으로 책임진다. 그러나 OMM이 실제로 어떻게 거래·헷지·수익화하는지에 대해서는 직접 증거가 거의 없었다."},

            {'kind':'h','text':"데이터 — 계좌 단위(account-level)의 KOSPI 200 데이터"},
            {'kind':'p','text':"한국거래소(KRX)로부터 익명화된 계좌 ID 기반의 옵션·선물 전 거래 기록을 확보했다. 2010년 1월 ~ 2014년 6월, 1,610,10개 계좌, 10억 건 이상의 옵션 거래를 포함한다. 두 시장이 동일 거래소에서 운영되어 OMM의 옵션 포지션과 기초자산(선물) 헷지 포지션을 동시에 관측할 수 있다."},

            {'kind':'h','text':"네 가지 주요 발견"},
            {'kind':'li','num':'1.','label':'High Profitability','text':"OMM은 평균 월 521백만 원의 옵션 PNL과 5.8–11.3의 샤프비율을 달성, 영업일의 74%에서 양(+)의 PNL."},
            {'kind':'li','num':'2.','label':'Rare Delta Hedging','text':"43명 중 단 4명(9.3%)만이 일관되게 델타헷지를 수행. 통념과 정면으로 배치되는 결과."},
            {'kind':'li','num':'3.','label':'Rapid Inventory Reversion','text':"재고의 38–48%가 5분 이내에 역전. 주된 도구는 옵션 지정가 주문(passive limit order)."},
            {'kind':'li','num':'4.','label':'Generalization to S&P 500','text':"S&P 500 OMM도 일평균 거래량이 순포지션 변화의 32배. 빠른 재고 회전 패턴 동일."},
        ],
    )

    # =========================================================================
    # 2. Data & Market Structure
    # =========================================================================
    build_content(
        prs,
        title="2. Data & Market Structure",
        subtitle="2.1 KOSPI 200 옵션·선물 시장의 특수성",
        blocks=[
            {'kind':'p','text':"KOSPI 200은 한국 200대 대형주를 추종하는 대표지수이며, 표본기간 동안 세계 최대 규모의 옵션시장이었다. 2010~2014 사이 옵션 86억 계약, 명목 거래대금 약 1.32경 원(약 1.2조 USD)이 체결됐다. WFE 기준 2011년 한국거래소가 글로벌 주가지수 옵션 거래량 1위였고, 이후 3개 시장(인도 NSE·EUREX·Cboe)을 합친 것보다 많았다."},

            {'kind':'h','text':"본 연구에 적합한 네 가지 시장 특성"},
            {'kind':'li','num':'1.','label':'World-leading Liquidity','text':"일평균 약 10억 USD 상당 옵션 거래, S&P 500 옵션과 비등한 규모. ATM·근ATM 집중도와 호가 스프레드 구조도 글로벌 시장과 유사."},
            {'kind':'li','num':'2.','label':'Modern Order-Driven Market','text':"중앙집중식 전자 지정가 주문장(LOB) 운영, 유동성 공급 인센티브 없음. 누구나 호가를 제출할 수 있어 OMM이 자기지정(self-appointed)으로 활동."},
            {'kind':'li','num':'3.','label':'Global OMM Participation','text':"Optiver, SIG, IMC, Virtu 등 글로벌 시장조성사 대부분 참여. 식별된 43명의 OMM 중 39명이 외국계."},
            {'kind':'li','num':'4.','label':'Single Hedging Asset','text':"옵션의 효율적 헷지가 KOSPI 200 선물에 한정 — 옵션·기초자산 시장 전체를 동시에 관측 가능."},

            {'kind':'h','text':"시장 구조의 두 가지 의미"},
            {'kind':'p','text':"동일 거래소·동일 수수료 구조 덕분에 옵션과 선물 사이의 헷지 거래에 인위적 마찰이 없다. 따라서 \"OMM이 헷지를 안 한다\"는 결과가 제도적 비용 때문이 아닌, 합리적 의사결정의 결과임을 의미한다."},
        ],
    )

    build_content(
        prs,
        title="2. Data & Market Structure",
        subtitle="2.2 OMM 식별 — 161,010개 계좌 → 43명의 시장조성자",
        blocks=[
            {'kind':'p','text':"KRX 데이터는 계좌가 OMM인지 별도 표시하지 않으므로, 거래·호가 패턴을 통해 추정한다. 시장조성자의 정의는 Harris(2002)를 따른다 — \"양방향 호가를 빈번히 갱신하고, 매수 후 매도/매도 후 매수를 반복해 큰 재고를 회피하며, 두 방향 균형을 유도하는 가격을 탐색하는 자\"."},

            {'kind':'h','text':"식별 절차 (단계별 필터)"},
            {'kind':'li','num':'STEP 1','label':'Top Passive Volume','text':"일별 패시브(유동성 공급) 거래대금 상위 30개 → 표본기간 동안 적어도 한 번 진입한 959개 계좌."},
            {'kind':'li','num':'STEP 2','label':'Two-Sided Quoting','text':"가장 활발한 종목 기준으로 거래일의 80% 이상 양방향 호가를 10일 이상 유지 → 90개 후보."},
            {'kind':'li','num':'STEP 3','label':'Volume Threshold','text':"월간 옵션 거래량 80,083 계약 초과(하위 5분위 제외) — 우연한 통과 계좌 제거."},
            {'kind':'li','num':'STEP 4','label':'BBO Competitiveness','text':"양방향 BBO 일치 시간 비율 14% 이상(상위 30%) — 실제 가격 발견에 기여."},
            {'kind':'li','num':'STEP 5','label':'Overnight Delta Limit','text':"평균 절대 델타 / 일평균 옵션거래량 < 0.0035 — 방향성 베팅 계좌 배제."},

            {'kind':'h','text':"식별 결과 — 표본의 통계적 특징"},
            {'kind':'p','text':"43개 OMM 계좌, 누적 294 계좌월. 동시에 활동하는 OMM은 평균 6명에 불과(높은 진입·퇴출률). 일평균 옵션 거래량 206,864 계약이지만 야간 보유 포지션은 단 4,431 계약 — 거래량 대비 재고가 극히 작은 전형적 시장조성 패턴. 양방향 호가 유지 시간 78%, 양방향 BBO 일치 32%로 매우 경쟁적이다."},
        ],
    )

    # =========================================================================
    # 3. Profitability
    # =========================================================================
    build_content(
        prs,
        title="3. Profitability of OMMs",
        subtitle="3.1 수익성 — \"매일 조금씩 그러나 꾸준히\"",
        blocks=[
            {'kind':'p','text':"OMM의 수익성을 측정하기 위해 일별 PNL을 각 계좌의 표본기간 최대 증거금으로 스케일한 일간 수익률을 사용한다(Baron et al. 2019의 방식). 이는 마진 기반 파생거래의 자본 효율성을 보다 정확히 반영한다."},

            {'kind':'h','text':"수익성의 핵심 지표"},
            {'kind':'li','num':'1.','label':'월 옵션 PNL','text':"평균 521.3 백만 원(약 47만 USD), 중앙값 285.2 백만 원, 25분위 74.8 백만 원으로 항상 양(+)의 영역."},
            {'kind':'li','num':'2.','label':'양(+)의 PNL 일수 비중','text':"월 평균 74% — 수의 트레이더(약 50%)보다 월등, 미국 주식 HFT(거의 100%)보다 낮은 수준."},
            {'kind':'li','num':'3.','label':'연환산 샤프비율','text':"옵션 단독 9.5, 옵션+선물 5.8 — 미국 주식 HFT와 견줄 만한 위험조정 성과."},
            {'kind':'li','num':'4.','label':'최악일 손실','text':"월 최악일 평균 -0.39%, 계좌 평균 최악일 -2.69% — 비시장조성 거래자(-28.5%)와 압도적 차이."},
            {'kind':'li','num':'5.','label':'스프레드 캡처율','text':"이론상 절반 스프레드 대비 평균 12% 캡처. 패시브 거래만 보면 43.6%까지 상승."},

            {'kind':'h','text':"비교 — 방향성 전략 대비"},
            {'kind':'p','text':"동일 시장에서 가장 성과가 좋은 방향성 전략(변동성 매도)은 월 PNL 41.5백만 원, 샤프비율 1.90에 그친다. 시장조성은 평균 12배 이상의 PNL과 약 5배 이상의 샤프비율을 달성, 위험조정 효율 면에서 명확히 우위에 있다."},
        ],
    )

    build_content(
        prs,
        title="3. Profitability of OMMs",
        subtitle="3.2 수익성의 결정요인 — 시장 조건 × 계좌 특성",
        blocks=[
            {'kind':'p','text':"OMM의 일별 수익률을 설명변수에 회귀분석한 결과, 시장 조건과 계좌 고유 특성이 모두 통계적으로 유의하게 작용한다. 마이크로스트럭처 이론의 핵심 예측과도 일치한다."},

            {'kind':'h','text':"시장 조건 (Column 1)"},
            {'kind':'li','num':'+','label':'Volatility','text':"변동성↑ → 수익성↑ (옵션 프리미엄 확대 효과)."},
            {'kind':'li','num':'+','label':'Trading Volume','text':"거래량↑ → 수익성↑ (유동성 수요 증가)."},
            {'kind':'li','num':'+','label':'Bid-Ask Spread','text':"스프레드↑ → 수익성↑ (단위 거래당 마진 확대)."},
            {'kind':'li','num':'·','label':'Index Return','text':"방향성 노출이 거의 없어 지수 수익률과의 상관성 무의미."},
            {'kind':'li','num':'·','label':'OMM Competition','text':"동시 활동 OMM 수는 단기 수익에 유의한 영향 없음."},

            {'kind':'h','text':"계좌 특성 (Column 2)"},
            {'kind':'li','num':'+','label':'Daily Volume','text':"거래량 ↑ → 수익률 ↑ (계수 0.010, t = 18.34) — 자본회전이 곧 수익."},
            {'kind':'li','num':'+','label':'Passive Trade Share','text':"패시브 비중 ↑ → 수익률 ↑ (0.070, t = 18.30) — 유동성 공급의 핵심 가치."},
            {'kind':'li','num':'+','label':'Foreign OMM Dummy','text':"외국계 OMM이 평균 +2.2% 우위 (t = 10.26) — 글로벌 운영 노하우 반영."},
            {'kind':'li','num':'−','label':'Account Time Trend','text':"진입 후 시간이 지날수록 성과 점진 하락 — 신규 경쟁자 진입 압력."},

            {'kind':'h','text':"진입·퇴출 패턴 (Section 3.4)"},
            {'kind':'p','text':"OMM은 진입 직후부터 일평균 약 30백만 원의 PNL을 안정적으로 시현한다. 퇴출은 급격한 손실이 아닌 점진적 수익성 악화의 결과 — 퇴출 10일 전부터 PNL이 평균 20백만 원 수준으로 둔화되며, 마지막 3일에는 평균 음(-)의 영역에 진입한다. 즉, 시장조성업은 \"붕괴\"가 아니라 \"경쟁 잠식\"으로 끝난다."},
        ],
    )

    # =========================================================================
    # 4. Delta Hedging
    # =========================================================================
    build_content(
        prs,
        title="4. Delta Hedging Behavior",
        subtitle="4.1 통념의 붕괴 — \"OMM은 거의 델타헷지를 하지 않는다\"",
        blocks=[
            {'kind':'p','text':"옵션 이론의 통념은 \"OMM은 항상 풀 델타헷지를 한다\"는 것이다. 본 연구는 계좌 단위로 옵션 델타와 선물 델타를 모두 관측해 이 가정을 직접 검증한다. 결과는 통념과 정면으로 충돌한다."},

            {'kind':'h','text':"43명 OMM의 헷지 행태 분포"},
            {'kind':'li','num':'·','label':'8명','text':"표본기간 동안 선물을 단 한 번도 거래하지 않음."},
            {'kind':'li','num':'·','label':'7명','text':"옵션·선물을 모두 거래하지만 일중(intraday)에 한정."},
            {'kind':'li','num':'·','label':'12명','text':"야간 옵션 포지션을 보유하면서도 선물로 헷지하지 않음."},
            {'kind':'li','num':'·','label':'18명','text':"옵션·선물을 모두 야간 보유 — 헷지 가능성이 있는 부분집합."},

            {'kind':'h','text':"헷지 비율 회귀분석 (18명 부분집합)"},
            {'kind':'p','text':"각 계좌별로 일말 선물 델타를 옵션 델타에 회귀해 헷지 비율(slope)을 추정한다. 풀 헷지라면 −1.0이 나와야 한다. 평균 슬로프 −0.2, 중앙값 −0.1로 절대값이 0.5에도 미치지 못한다. 평균 R² 21.8%로 일관성도 낮다."},

            {'kind':'h','text':"비모수 검정 — 의미 있는 헷지일의 비중"},
            {'kind':'p','text':"옵션·선물 델타가 반대 부호이고, 선물 델타가 옵션 델타의 0.5 ~ 1.5배 구간에 위치하는 \"의미 있는 헷지\" 일을 정의한다. 평균은 단 9.1%(범위 0% ~ 42.6%)로, 무작위 포지셔닝(약 50%)에도 못 미친다. 결국 43명 중 4명만 일관된 델타헷지 패턴을 보인다."},
        ],
    )

    build_content(
        prs,
        title="4. Delta Hedging Behavior",
        subtitle="4.2 누가, 언제 델타헷지를 하는가? — 4명의 델타헷저 분석",
        blocks=[
            {'kind':'p','text':"전체의 9.3%인 4명만이 일관된 델타헷지를 시행한다. 이들의 행태를 분석하면 \"왜 다수는 헷지하지 않는가\"에 대한 답을 거꾸로 읽어낼 수 있다. 이들의 헷지 결정은 비용·편익의 합리적 계산으로 설명된다."},

            {'kind':'h','text':"4명의 헷지 강도 — 풀헷지에 가까운 것은 단 1명"},
            {'kind':'li','num':'#1','label':'Limited Hedger','text':"슬로프 −0.27, R² 36% — 부호는 맞추되 일관성 낮음."},
            {'kind':'li','num':'#2','label':'Classical Hedger','text':"슬로프 −0.90, R² 96% — 교과서적인 풀 헷지에 근접."},
            {'kind':'li','num':'#3','label':'Moderate Hedger','text':"슬로프 −0.49, R² 52% — 절반 정도만 헷지."},
            {'kind':'li','num':'#4','label':'Two-Regime Hedger','text':"두 가지 레짐 사이를 전환하는 적응형 헷지."},

            {'kind':'h','text':"헷지의 결정요인 — 비용/편익 모형과 정합적"},
            {'kind':'li','num':'A.','label':'Holding Period','text':"델타헷저의 일중 평균 보유시간 46.1분 vs. 비헷저 27.5분 — 보유가 긴 OMM이 헷지로부터의 위험감축 효익이 큼."},
            {'kind':'li','num':'B.','label':'Contract Liquidity','text':"델타헷저는 비헷저보다 덜 유동적인 종목에서 거래(평균 거래량 391K vs. 648K) — 옵션 시장 내 재고 회전이 어려운 환경에서 헷지가 보완재로 기능."},
            {'kind':'li','num':'C.','label':'Intraday Pattern','text':"헷지 비율은 개장 직후 −0.291, 정오 −0.205, 마감 −0.317의 역U자(inverse-U)패턴 — 야간 위험 직전 강화."},
            {'kind':'li','num':'D.','label':'Stress Response','text':"변동성 급등·재고 충격 시 헷지 강도 거의 두 배. 즉, 델타헷지는 \"기본 도구\"가 아닌 \"스트레스 대응\" 도구."},

            {'kind':'h','text':"실행 — 패시브 우위, 상대적 비용 최소화"},
            {'kind':'p','text':"4명의 델타헷저도 선물 거래의 절반은 지정가(passive)로 체결시켜 비용을 통제한다. 즉, 선택적·간헐적 풀헷지조차도 \"비용 최소화 + 옵션 재고관리\"라는 시장조성 원칙에 종속된다."},
        ],
    )

    # =========================================================================
    # 5. Active Inventory Management
    # =========================================================================
    build_content(
        prs,
        title="5. Active Inventory Management",
        subtitle="5.1 핵심 메커니즘 — \"분 단위 재고 회전\"",
        blocks=[
            {'kind':'p','text':"OMM의 진짜 위험관리 도구는 헷지가 아니라 능동적 재고관리(active inventory management)다. Amihud and Mendelson(1980) 모형과 정합적으로, OMM은 재고가 한쪽으로 쏠리면 호가를 비대칭적으로 움직여 반대 방향 주문을 유도한다. 이 과정에서 재고는 자연스럽게 평균회귀한다."},

            {'kind':'h','text':"5분 간격 재고 변화의 자기회귀 검정"},
            {'kind':'li','num':'·','label':'AR(1) only','text':"평균 계수 −0.383, t = −36.7 — 재고 변화의 약 38%가 다음 5분 내 역전."},
            {'kind':'li','num':'·','label':'AR(1) + AR(2)','text':"AR(1) −0.479, AR(2) −0.219 → 합 약 −0.7 — 약 70%의 불균형이 10분 이내 상쇄."},
            {'kind':'li','num':'·','label':'Robustness','text':"동시·시차 지수 수익률을 통제해도 결과 변동 없음."},

            {'kind':'h','text':"분해 — 어떤 거래로 재고를 줄이는가?"},
            {'kind':'li','num':'1.','label':'Passive Options Trades','text':"−28.3% (가장 큰 기여) — 호가 조정으로 자연주문 흐름을 끌어들임. 시장조성의 본질."},
            {'kind':'li','num':'2.','label':'Aggressive Options Trades','text':"−9.5% — 임계 수준에 도달했을 때만 시장가 주문(taker) 사용."},
            {'kind':'li','num':'3.','label':'Futures Trades','text':"−0.5% — 사실상 무시 가능 수준. 헷지가 아닌 옵션 시장 자체에서 위험관리."},

            {'kind':'h','text':"이론적 함의 — 마이크로스트럭처가 옳다"},
            {'kind':'p','text':"빠른 재고 회전이 가능한 유동시장에서는 헷지의 거래비용(스프레드, 마켓임팩트)이 위험감축 편익을 상쇄한다. 따라서 \"OMM = 풀 델타헷저\"라는 통념은 모든 시장에 보편적이지 않다. KOSPI 200처럼 매우 유동적인 옵션시장에서는 마이크로스트럭처 이론이 옵션 이론보다 행동을 더 잘 설명한다."},
        ],
    )

    # =========================================================================
    # 6. S&P 500 Evidence
    # =========================================================================
    build_content(
        prs,
        title="6. Evidence from S&P 500 Options",
        subtitle="6.1 \"빠른 재고 회전\"은 미국 시장에서도 작동한다",
        blocks=[
            {'kind':'p','text':"KOSPI 결과가 한국 시장에 국한되는지 확인하기 위해, 1996~2018년 23년치 Cboe Open-Close Volume 데이터로 S&P 500 지수옵션을 분석한다. 계좌 단위 데이터는 없지만, OMM 그룹의 일별 매수/매도·개시/청산 거래량을 모두 관측 가능 — S&P 500 옵션이 Cboe에서만 거래되므로 OMM 활동 전수 포착이 가능하다."},

            {'kind':'h','text':"가설 — \"빠른 회전\"의 흔적은 거래량과 순포지션 변화의 비율에 남는다"},
            {'kind':'p','text':"OMM이 모든 옵션 거래를 즉시 풀 델타헷지한다는 통념대로라면, 헷지 수요는 일평균 거래량과 비례해야 한다. 반면 OMM이 같은 종목 내·유사 종목 사이로 재고를 빠르게 회전시킨다면, 일말 순포지션 변화는 거래량보다 훨씬 작아야 한다."},

            {'kind':'h','text':"핵심 통계 — 거래량은 순변화의 32배"},
            {'kind':'li','num':'·','label':'Daily OMM Volume','text':"평균 415,840 계약 / 중앙값 370,189 계약. ATM·근ATM 집중, 시간이 지날수록 증가."},
            {'kind':'li','num':'·','label':'Naïve Hedge Need','text':"모든 거래를 즉시 헷지한다고 가정 시 일평균 110,953 계약 델타. 명목 약 158억 USD."},
            {'kind':'li','num':'·','label':'Actual Net Δ Position Change','text':"평균 3,466 계약 / 중앙값 2,260 계약 — 실제 헷지가 필요할 \"잔여\" 노출."},
            {'kind':'li','num':'·','label':'Volume / Net Change','text':"평균 32배 / 중앙값 40배 — KOSPI 결과와 동일한 \"빠른 회전\" 패턴."},

            {'kind':'h','text':"광범위한 함의 — 0DTE 시장 안정성 논쟁에 대한 시사점"},
            {'kind':'p','text':"최근 0DTE 옵션의 폭발적 거래량이 OMM의 헷지 수요를 통해 기초자산 변동성을 키운다는 우려가 있다. 그러나 본 연구는 OMM이 거래량의 대부분을 빠른 재고 회전으로 흡수한다는 사실을 보여준다. 따라서 0DTE 옵션의 \"실효 헷지 수요\"는 명목 거래량이 시사하는 것보다 훨씬 작을 수 있다."},
        ],
    )

    # =========================================================================
    # 7. Conclusion
    # =========================================================================
    build_content(
        prs,
        title="7. Conclusion & Implications",
        subtitle="7.1 발견의 종합과 향후 연구 방향",
        blocks=[
            {'kind':'p','text':"본 연구는 계좌 단위 KOSPI 200 데이터와 23년치 S&P 500 데이터를 결합해 옵션 시장조성자(OMM)의 거래·헷지·수익 구조를 직접 분석한 첫 시도다. 핵심 발견은 \"OMM은 헷지가 아니라 분 단위 재고 회전으로 위험을 관리한다\"는 것이다."},

            {'kind':'h','text':"네 가지 핵심 결론"},
            {'kind':'li','num':'1.','label':'Profit Stability','text':"시장조성은 매우 수익성이 높고(샤프 5.8–11.3) 손실일도 제한적이다. 퇴출은 점진적 경쟁 잠식으로 발생한다."},
            {'kind':'li','num':'2.','label':'Delta Hedging Is Selective','text':"43명 중 4명만 일관된 델타헷지. 헷저는 보유가 길고 비유동적 종목을 다루는 OMM에 집중된다."},
            {'kind':'li','num':'3.','label':'Inventory Management Dominates','text':"재고 불균형의 ~70%가 10분 내 해소. 주된 도구는 패시브 옵션 호가 조정(약 28%)."},
            {'kind':'li','num':'4.','label':'Generalizable to S&P 500','text':"OMM 거래량이 순포지션 변화의 32배로, 빠른 회전이 미국 시장의 1차 위험관리 메커니즘."},

            {'kind':'h','text':"이론·실무에 대한 함의"},
            {'kind':'li','num':'A.','label':'Theoretical','text':"마이크로스트럭처 이론이 유동성 높은 옵션시장의 행동을 더 잘 설명. 옵션 이론의 \"풀 델타헷지\" 가정은 보편이 아닌 시장 조건의 함수."},
            {'kind':'li','num':'B.','label':'Market Stability','text':"0DTE 등 명목 거래량 급증이 곧 헷지 수요 급증을 의미하지 않는다 — 시장 영향 평가 시 \"순 포지션 변화\"를 봐야 한다."},
            {'kind':'li','num':'C.','label':'Future Research','text':"비유동적 미국 주식옵션 등 \"긴 보유 + 좁은 옵션 시장\" 환경에서는 델타헷지 의존도가 클 가능성. 계좌 단위 데이터 확보가 후속 과제."},

            {'kind':'h','text':"한 줄 요약"},
            {'kind':'p','text':"옵션 시장조성자의 진짜 무기는 \"기초자산을 사는 손\"이 아니라 \"호가를 움직이는 손\"이다."},
        ],
    )

    # ---- Save ----
    out = os.path.join(ROOT, 'Y-PJ-FoRM_KOSPI200_DeltaHedging.pptx')
    prs.save(out)
    print(f"saved -> {out}")
    print(f"slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
